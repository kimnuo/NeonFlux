#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
NeonFlux Final Presentation Generator
Generates sample.pptx with actual implementation details
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette (NeonFlux theme) ──
BG_DARK      = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD      = RGBColor(0x16, 0x21, 0x3E)
NEON_CYAN    = RGBColor(0x00, 0xFF, 0xFF)
NEON_MAGENTA = RGBColor(0xFF, 0x00, 0xFF)
NEON_YELLOW  = RGBColor(0xFF, 0xFF, 0x00)
NEON_GREEN   = RGBColor(0x00, 0xFF, 0x66)
NEON_RED     = RGBColor(0xFF, 0x44, 0x44)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY     = RGBColor(0x99, 0x99, 0x99)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Malgun Gothic'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=LIGHT_GRAY, bullet_color=NEON_CYAN, font_name='Malgun Gothic'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.level = 0
    return txBox

# ══════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, BG_DARK)

# Top accent line
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

# Project title
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
             "NeonFlux", 72, NEON_CYAN, True, PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.8),
             "네온 시티를 질주하는 드리프트 레이싱 게임", 32, NEON_MAGENTA, False, PP_ALIGN.CENTER)

# Divider
add_shape(slide, Inches(4.5), Inches(3.9), Inches(4.3), Inches(0.04), NEON_YELLOW)

# Info
add_text_box(slide, Inches(2), Inches(4.3), Inches(9.3), Inches(0.5),
             "최종 발표 자료  |  2025년 6월", 24, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.5),
             "김택준  |  게임소프트웨어 26001", 24, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(5.6), Inches(9.3), Inches(0.5),
             "Unity 2022.3.62f3 (URP)  |  Android 타겟", 20, MID_GRAY, False, PP_ALIGN.CENTER)

# Bottom accent line
add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 2: 개발 개요
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "1. 개발 개요", 36, NEON_CYAN, True)

# Left card - Project Info
card = add_shape(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5), BG_CARD, NEON_CYAN, 1.5)
add_text_box(slide, Inches(0.9), Inches(1.4), Inches(5.2), Inches(0.5),
             "프로젝트 정보", 24, NEON_CYAN, True)

items_left = [
    "▸ 프로젝트명: NeonFlux (네온플럭스)",
    "▸ 장르: 드리프트 레이싱 (수동 조작)",
    "▸ 타겟 플랫폼: Android (세로 모드)",
    "▸ 개발 기간: 2025년 4월 ~ 6월",
    "▸ 엔진: Unity 2022.3.62f3 (URP)",
    "▸ 개발 인원: 1인 개발 (김택준)",
    "▸ AI 도구 활용: Qwen Code (LLM 기반 코드 생성/디버깅)",
]
add_bullet_list(slide, Inches(0.9), Inches(2.0), Inches(5.2), Inches(4.5), items_left, 18, LIGHT_GRAY)

# Right card - Key Features
card = add_shape(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5), BG_CARD, NEON_MAGENTA, 1.5)
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.2), Inches(0.5),
             "핵심 기능", 24, NEON_MAGENTA, True)

items_right = [
    "▸ 물리 기반 차량 조향 및 드리프트 시스템",
    "▸ 지형 추종 (Ground Following) + Pitch 정렬",
    "▸ 주행 거리 기반 실시간 점수 누적",
    "▸ 게임 기록 영구 저장 및 리더보드 (상위 10개)",
    "▸ 스테이지 클리어 & 게임 오버 UI",
    "▸ AutoUIBuilder - 런타임 UI 자동 생성/연결",
    "▸ 카메라 추종 + 드리프트 스모크 파티클",
]
add_bullet_list(slide, Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.5), items_right, 18, LIGHT_GRAY)

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 3: 개발 환경
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "2. 개발 환경", 36, NEON_CYAN, True)

card = add_shape(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), BG_CARD, NEON_CYAN, 1.5)

# Table-like layout
env_data = [
    ("항목", "내용"),
    ("Unity 버전", "2022.3.62f3 (URP)"),
    ("스크립팅 백엔드", "IL2CPP"),
    ("타겟 플랫폼", "Android (ARM64 + ARMv7, SDK 34)"),
    ("해상도", "세로 모드 1080x1920 기준"),
    ("조작 방식", "FixedSlideBar UI + Touch/Mouse 드래그"),
    ("저장 방식", "JSON 직렬화 (Application.persistentDataPath)"),
    ("에셋", "EasyRoads3D v3, AsserStore 에셋 응용"),
    ("버전 관리", "Git + GitHub Copilot instructions"),
    ("AI 도구", "Qwen Code (LLM) - 코드 생성, 리팩토링, 디버깅"),
]

y_pos = 1.5
for i, (key, val) in enumerate(env_data):
    is_header = (i == 0)
    bg_c = RGBColor(0x0D, 0x1B, 0x33) if is_header else (RGBColor(0x12, 0x25, 0x44) if i % 2 == 0 else RGBColor(0x0F, 0x1F, 0x3A))
    txt_c = NEON_CYAN if is_header else LIGHT_GRAY
    bld = is_header

    add_shape(slide, Inches(0.8), Inches(y_pos), Inches(3.5), Inches(0.48), bg_c)
    add_shape(slide, Inches(4.3), Inches(y_pos), Inches(8.1), Inches(0.48), bg_c)
    add_text_box(slide, Inches(0.9), Inches(y_pos + 0.04), Inches(3.3), Inches(0.4),
                 key, 18, txt_c, bld)
    add_text_box(slide, Inches(4.4), Inches(y_pos + 0.04), Inches(7.9), Inches(0.4),
                 val, 18, txt_c, bld)
    y_pos += 0.52

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 4: 시스템 아키텍처
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "3. 시스템 아키텍처", 36, NEON_CYAN, True)

# Architecture cards
cards = [
    ("Architecture (상태 관리)", NEON_CYAN, [
        "GameManager: GameState FSM (MainMenu/Playing/StageClear/GameOver)",
        "Singleton<T>: 전역 매니저 기반 구조",
        "InputManager: SlideBar/Touch/Mouse 통합 입력 처리",
    ]),
    ("Entities (게임 객체)", NEON_MAGENTA, [
        "PlayerController: 차량 물리, 조향, 드리프트, 점수 (~1366줄)",
        "SpeedBoostPad: 속도 증가 트리거 (ResetForReplay 지원)",
        "ScorePickup: 점수 수집 아이템",
        "EngineAudioController: 엔진/스키드 오디오",
    ]),
    ("Managers (시스템)", NEON_YELLOW, [
        "SaveManager: JSON 직렬화, 리더보드, 최고 점수",
        "LevelDifficultyManager: 시간 기반 장애물 밀도 계산",
    ]),
    ("UI (인터페이스)", NEON_GREEN, [
        "AutoUIBuilder: 런타임 UI 자동 생성/연결 (~584줄)",
        "GameOverUI / StageClearUI: 상태별 결과 화면",
        "LeaderboardUI: 상위 10개 기록 표시",
        "FixedSlideBar: 조향용 UI 슬라이더",
    ]),
]

for i, (title, color, items) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = 0.6 + col * 6.2
    y = 1.3 + row * 2.9

    card = add_shape(slide, Inches(x), Inches(y), Inches(5.9), Inches(2.7), BG_CARD, color, 1.5)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 0.1), Inches(5.5), Inches(0.4),
                 title, 20, color, True)
    add_bullet_list(slide, Inches(x + 0.2), Inches(y + 0.55), Inches(5.5), Inches(2.0),
                    items, 14, LIGHT_GRAY)

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 5: 핵심 구현 - 차량 제어 시스템
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "4. 핵심 구현 - 차량 제어 시스템", 36, NEON_CYAN, True)

card = add_shape(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), BG_CARD, NEON_CYAN, 1.5)

items = [
    "▸ 물리 기반 구동: Rigidbody.velocity 기반 자동 가속 (cruiseSpeed → maxForwardSpeed)",
    "▸ Yaw 기반 조향: 속도 의존적 조향각 감소 (steerBySpeed 커브) + 헤딩 제한 (110도)",
    "▸ 드리프트 시스템:",
    "    - 최대 조향 입력 시 자동 드리프트 진입",
    "    - driftLateralGrip vs normalLateralGrip로 접지력 차별화",
    "    - 드리프트 중 forwardBoost로 속도 증가 (체감 보상)",
    "▸ 지형 추종 (Ground Following): 3개 레이캐스트 (중앙/좌/우)로 노멀 감지 + 보간",
    "▸ Pitch 정렬: 전방/후방 레이캐스트로 도로 경사각 반영",
    "▸ Step Assist: 작은 장애물 등반 지원 (선택적)",
    "▸ 휠 비주얼: 자식 휠 오브젝트 회전 + 조향각 시각화",
    "▸ 드리프트 스모크 파티클: 좌/우/중앙 이미터, 슬립량 기반 강도 조절",
    "▸ Airborne Clamp: 공중에서의 비정상 상승 방지",
    "▸ Front Impact Pitch Lock: 정면 충돌 후 0.2초간 상향 속도 잠금",
]
add_bullet_list(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.2), items, 16, LIGHT_GRAY)

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 6: 핵심 구현 - 점수 & 저장 시스템
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "5. 핵심 구현 - 점수 & 저장 시스템", 36, NEON_CYAN, True)

# Left card - Scoring
card = add_shape(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5), BG_CARD, NEON_GREEN, 1.5)
add_text_box(slide, Inches(0.9), Inches(1.4), Inches(5.2), Inches(0.5),
             "실시간 점수 시스템", 24, NEON_GREEN, True)

items_left = [
    "▸ 주행 거리 기반 자동 누적:",
    "    - Z축 (전진): 5m당 2점",
    "    - X축 (횡이동): 1m당 1점",
    "▸ ScorePickup: 충돌 시 보너스 점수 (10~50)",
    "▸ SpeedBoostPad: 충돌 시 maxForwardSpeed 영구 증가",
    "▸ GameManager.SetCurrentScore(): 실시간 동기화",
    "▸ 게임 오버 시 정확한 최종 점수 표시",
    "▸ 디지털 속도계: 실시간 km/h 표시 (우측 상단)",
]
add_bullet_list(slide, Inches(0.9), Inches(2.0), Inches(5.2), Inches(4.5), items_left, 16, LIGHT_GRAY)

# Right card - Save/Leaderboard
card = add_shape(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5), BG_CARD, NEON_YELLOW, 1.5)
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.2), Inches(0.5),
             "영구 저장 & 리더보드", 24, NEON_YELLOW, True)

items_right = [
    "▸ JSON 직렬화 (JsonUtility):",
    "    - Application.persistentDataPath/neonflux_user_data.json",
    "▸ 저장 데이터:",
    "    - 최고 점수, 현재 스테이지 레벨",
    "    - scoreHistory (최대 10개 기록)",
    "▸ LeaderboardEntry: 점수/스테이지/날짜/사유",
    "▸ 리더보드 UI: 상위 10개 순위 표시",
    "    - Gold/Silver/Bronze 컬러 (1~3위)",
    "    - MM/dd HH:mm 형식 날짜 표시",
    "▸ 게임 오버 화면에 기록 점수 표시",
    "▸ StageClearUI: 스테이지 클리어 시 긍정적 메시지",
]
add_bullet_list(slide, Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.5), items_right, 16, LIGHT_GRAY)

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 7: 핵심 구현 - UI 시스템
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "6. 핵심 구현 - UI 시스템", 36, NEON_CYAN, True)

card = add_shape(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), BG_CARD, NEON_CYAN, 1.5)

items = [
    "▸ AutoUIBuilder (런타임 UI 자동 생성/연결):",
    "    - Canvas, EventSystem 자동 생성",
    "    - MainMenu / Leaderboard / GameOver / StageClear 패널 자동 빌드",
    "    - 모든 UI 컴포넌트 참조 자동 와이어링 (리플렉션 활용)",
    "    - GameManager의 stageClearUI, gameOverUI 필드 자동 할당",
    "",
    "▸ GameState 기반 UI 전환 (GameManager.ApplyStateVisuals()):",
    "    - MainMenu: 시작 버튼 + 리더보드 토글 + 최고 점수",
    "    - Playing: 디지털 속도계 + 점수 HUD 표시",
    "    - GameOver: 점수 / 사유 / 최고 점수 / 기록 점수 / 버튼",
    "    - StageClear: 점수 / 긍정 메시지 / 최고 점수 / 기록 점수 / 다음 단계 버튼",
    "",
    "▸ SafeAreaController: 노치/제스처 바 대응 (모바일)",
    "▸ FixedSlideBar: 조향용 UI 슬라이더 (IPointerDownHandler/IDragHandler/IPointerUpHandler)",
    "▸ CameraFollow: LateUpdate 기반 부드러운 카메라 추종 (GameState.Playing 중만)",
]
add_bullet_list(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.2), items, 16, LIGHT_GRAY)

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 8: 에셋 활용 및 응용
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "7. 에셋 활용 및 응용", 36, NEON_CYAN, True)

card = add_shape(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), BG_CARD, NEON_MAGENTA, 1.5)

items = [
    "▸ EasyRoads3D v3:",
    "    - 도로 생성 및 커브/정션 시스템 활용",
    "    - RoadManager 스텁과 연동하여 장애물 스폰 위치 결정",
    "",
    "▸ AsserStore 에셋:",
    "    - 차량 모델, 파티클, 사운드 등 기본 에셋 응용",
    "    - 드리프트 스모크 파티클: 기존 에셋을 슬라이드/조향 입력에 연동",
    "    - 엔진 오디오: 속도/드리프트 상태에 따라 클립 전환 (maxRpmClip)",
    "",
    "▸ Google Ads (AdMob) 완전 제거:",
    "    - 초기 기획에 포함되었으나 빌드 충돌 문제로 전량 삭제",
    "    - GoogleMobileAdsPlugin, 관련 aar, iOS 플러그인, 의존성 패키지 전체 제거",
    "    - ExternalDependencyManager는 다른 플러그인 호환성 유지",
    "",
    "▸ Android 빌드 문제 해결:",
    "    - Kotlin/kotlinx-coroutines 버전 충돌 → 의존성 정리",
    "    - Target SDK 34, ARM64+ARMv7 듀얼 아키텍처 설정",
    "    - Custom Gradle Template 활성화",
]
add_bullet_list(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.2), items, 16, LIGHT_GRAY)

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 9: AI/LLM 활용
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "8. AI/LLM 활용 (Qwen Code)", 36, NEON_CYAN, True)

card = add_shape(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), BG_CARD, NEON_GREEN, 1.5)

items = [
    "▸ Qwen Code (LLM 기반 AI 코딩 도구) 활용:",
    "",
    "  1) 코드 생성 및 구현:",
    "     - GameOverUI, StageClearUI, LeaderboardUI, MainMenuUI 컴포넌트 생성",
    "     - AutoUIBuilder 런타임 UI 생성 시스템 구현 (~584줄)",
    "     - GameManager 상태 머신 및 UI 오케스트레이션",
    "     - SaveManager JSON 직렬화 및 리더보드 시스템",
    "",
    "  2) 리팩토링:",
    "     - PlayerController 중복 코드 제거 (OutOfBoundsPanel → GameOverPanel 통합)",
    "     - StageClearUI를 GameOverUI 기반으로 재구성 (긍정적 메시지)",
    "     - GameManager에 StageClear 상태 처리 추가",
    "",
    "  3) 디버깅 및 문제 해결:",
    "     - Android Gradle 빌드 실패 분석 및 의존성 충돌 해결",
    "     - Google Ads 완전 제거 및 빌드 정리",
    "     - 컴파일 오류 (CS1061, CS0103) 추적 및 수정",
    "",
    "  4) 문서화:",
    "     - 작업 내역서 작성",
    "     - Readme.txt 실행 가이드 작성",
    "     - 최종 발표 자료 PPT 생성",
]
add_bullet_list(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.2), items, 16, LIGHT_GRAY)

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 10: 데이터 흐름
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "9. 데이터 흐름 다이어그램", 36, NEON_CYAN, True)

card = add_shape(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), BG_CARD, NEON_CYAN, 1.5)

items = [
    "▸ 게임 진행 중 (실시간):",
    "    PlayerController → UpdateScoreTracking() → GameManager.SetCurrentScore(score)",
    "",
    "▸ 게임 종료 시 (장애물 충돌 / 코스 이탈):",
    "    GameManager.SetGameOver(reason) → RecordGameScore(reason)",
    "    → SaveManager.RecordGameScore(score, reason)",
    "      1. CurrentData.highScore 갱신 (현재 > 기존 최고)",
    "      2. LeaderboardEntry 생성 (score, stageLevel, date, reason)",
    "      3. scoreHistory.Add(entry) → 10개 초과 시 오래된 항목 제거",
    "      4. File.WriteAllText()로 디스크 저장",
    "    → ApplyStateVisuals() → GameOverUI.ShowGameOver(score, reason)",
    "",
    "▸ 스테이지 클리어 시:",
    "    GameManager.CompleteStage() → RecordGameScore(\"clear\")",
    "    → SaveManager.RecordGameScore(score, \"clear\")",
    "    → currentStageLevel++ → SaveData()",
    "    → ApplyStateVisuals() → StageClearUI.ShowStageClear(score, reason)",
    "",
    "▸ 메인 메뉴:",
    "    GameManager.GoToMainMenu() → mainMenuUI.SetActive(true)",
    "    → MainMenuUI.OnEnable() → UpdateHighScoreDisplay()",
]
add_bullet_list(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.2), items, 15, LIGHT_GRAY)

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 11: 트러블슈팅
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "10. 트러블슈팅", 36, NEON_CYAN, True)

card = add_shape(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), BG_CARD, NEON_RED, 1.5)

items = [
    "▸ Android Gradle 빌드 실패 (CommandInvokationFailure):",
    "    - kotlin-stdlib 2.1.0 dexing 실패 → 의존성 패키지 제거",
    "    - kotlinx-coroutines 1.8.0 dexing 실패 → 동일 처리",
    "    - AndroidResolverDependencies.xml ABI가 armeabi-v7a만 설정 → arm64-v8a 추가",
    "",
    "▸ Google Ads (AdMob) 충돌:",
    "    - GoogleMobileAdsPlugin.androidlib, googlemobileads-unity.aar 등 전체 제거",
    "    - mainTemplate.gradle에서 의존성 자동 정리",
    "    - ExternalDependencyManager 유지 (다른 플러그인 호환)",
    "",
    "▸ Unity 컴파일 오류:",
    "    - CS1061: 'GameOverUI'에 retryButton 정의 없음 → lastScoreText로 변경",
    "    - CS0103: '_isOutOfBounds' 미존재 → 남은 참조 제거",
    "    - PlayerController의 StageClear 관련 코드 제거 → StageClearUI로 분리",
    "",
    "▸ Gradle Template 문제:",
    "    - android.suppressUnsupportedCompileSdk=36 제거",
    "    - android.jetifier.ignorelist 제거",
    "    - Custom Gradle Template 활성화 (useCustomMainGradleTemplate, useCustomGradleSettingsTemplate)",
]
add_bullet_list(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.2), items, 15, LIGHT_GRAY)

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 12: 구현 현황
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "11. 구현 현황", 36, NEON_CYAN, True)

# Left - Implemented
card = add_shape(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5), BG_CARD, NEON_GREEN, 1.5)
add_text_box(slide, Inches(0.9), Inches(1.4), Inches(5.2), Inches(0.5),
             "완료 (구현됨)", 24, NEON_GREEN, True)

items_left = [
    "✅ 차량 물리 기반 조향 시스템",
    "✅ 드리프트 시스템 (자동 진입, 속도 보상)",
    "✅ 지형 추종 + Pitch 정렬",
    "✅ 주행 거리 기반 실시간 점수",
    "✅ ScorePickup / SpeedBoostPad",
    "✅ 게임 기록 영구 저장 (JSON)",
    "✅ 리더보드 (상위 10개)",
    "✅ GameOverUI / StageClearUI",
    "✅ AutoUIBuilder (런타임 UI 생성)",
    "✅ 카메라 추종",
    "✅ 드리프트 스모크 파티클",
    "✅ 엔진/스키드 오디오",
    "✅ Android 빌드 (SDK 34, ARM64+ARMv7)",
]
add_bullet_list(slide, Inches(0.9), Inches(2.0), Inches(5.2), Inches(4.5), items_left, 16, LIGHT_GRAY)

# Right - Future
card = add_shape(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5), BG_CARD, NEON_YELLOW, 1.5)
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.2), Inches(0.5),
             "향후 계획 (미구현)", 24, NEON_YELLOW, True)

items_right = [
    "⬜ UpgradeStatsSO 레벨별 데이터 확장",
    "⬜ 실제 도로 생성 시스템 (RoadManager 구현)",
    "⬜ 장애물 스폰 시스템",
    "⬜ 업그레이드/재화 시스템",
    "⬜ 단위 테스트 도입",
    "⬜ 추가 스테이지/코스 디자인",
    "⬜ 멀티터치 조작 지원",
    "⬜ achievements/랭킹 시스템",
]
add_bullet_list(slide, Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.5), items_right, 18, LIGHT_GRAY)

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 13: 최종 결과물
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "12. 최종 결과물", 36, NEON_CYAN, True)

card = add_shape(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), BG_CARD, NEON_CYAN, 1.5)

items = [
    "▸ 제출 목록:",
    "    1. 기획안.pdf",
    "    2. 최종 발표자료 (본 PPT)",
    "    3. 게임 실행 영상 파일 (5분 이상)",
    "    4. Readme.txt (Unity 버전 + 실행 참고 사항)",
    "    5. 소스 파일 전체 (Unity 프로젝트)",
    "",
    "▸ 실행 방법:",
    "    1. Unity Hub에서 Unity 2022.3.62f3 설치",
    "    2. Unity Hub → Projects → Add → 프로젝트 폴더 선택",
    "    3. 프로젝트 열기 후 자동 패키지 리졸브 대기",
    "    4. 에디터 재생 버튼으로 테스트",
    "       - PC: 좌우 화살표/A,D 키로 조향",
    "       - Android: File → Build Settings → Android → Build",
    "",
    "▸ Unity 버전: 2022.3.62f3 (URP)",
    "▸ 빌드 타겟: Android (ARM64 + ARMv7, SDK 34)",
    "▸ 스크립팅 백엔드: IL2CPP 권장",
]
add_bullet_list(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.2), items, 16, LIGHT_GRAY)

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ══════════════════════════════════════════════════════════
# SLIDE 14: Thank You
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), NEON_CYAN)

add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.2),
             "감사합니다", 60, NEON_CYAN, True, PP_ALIGN.CENTER)

add_shape(slide, Inches(4.5), Inches(3.5), Inches(4.3), Inches(0.04), NEON_YELLOW)

add_text_box(slide, Inches(2), Inches(4.0), Inches(9.3), Inches(0.6),
             "NeonFlux - 네온 시티를 질주하는 드리프트 레이싱", 28, NEON_MAGENTA, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(0.5),
             "김택준  |  게임소프트웨어 26001", 24, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.5),
             "Unity 2022.3.62f3  |  Android  |  1인 개발", 20, MID_GRAY, False, PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), NEON_MAGENTA)


# ── Save ──
output_path = r"D:\toy\Unity\NeonFlux\sample.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
